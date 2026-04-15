"""PPT generator using Builder/Factory pattern for slide creation.

Uses the BattleCardFlashTmp.pptx template:
  - Slide 1: Title page (kept as-is from template)
  - Slide 2: Content template (cloned for each content section)
  - Last slide: Thanks page (kept as-is from template)
"""
import copy
import os
from abc import ABC, abstractmethod
from datetime import date
from typing import List

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from models.entities import ComparisonResult

# Template path
TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "BattleCardFlashTmp.pptx")

# Color palette
COLOR_PRIMARY = RGBColor(0x16, 0x21, 0x3E)
COLOR_ACCENT = RGBColor(0x0F, 0x3D, 0x6E)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)

# Content area for 10" x 5.62" template slides
CONTENT_LEFT = Inches(0.4)
CONTENT_WIDTH = Inches(6.5)


def _clone_slide(prs: Presentation, template_slide) -> object:
    """Clone a slide (shapes + layout) and append it to the presentation."""
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default placeholder shapes from the new slide
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Deep-copy all shapes from the template slide
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    return new_slide


class SlideBuilder(ABC):
    """Abstract builder for PPT slides."""

    @abstractmethod
    def build(self, slide, result: ComparisonResult) -> None:
        ...


class FeatureTableSlideBuilder(SlideBuilder):
    """Builds the feature comparison table slide."""

    def build(self, slide, result: ComparisonResult) -> None:
        # Title
        txBox = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.25), CONTENT_WIDTH, Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Feature Comparison: {result.product_a.name} vs {result.product_b.name}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        # Table
        rows = len(result.feature_comparisons) + 1
        cols = 3
        left, top = CONTENT_LEFT, Inches(0.85)
        width = Inches(9.0)
        height = Inches(min(0.3 * rows, 4.5))
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Column widths
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(3.4)
        table.columns[2].width = Inches(3.4)

        # Header row
        headers = ["Feature", result.product_a.name, result.product_b.name]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            _style_cell(cell, Pt(10), bold=True, bg=COLOR_PRIMARY, fg=COLOR_WHITE)

        # Data rows
        for row_idx, (feature, val_a, val_b, weight) in enumerate(result.feature_comparisons, 1):
            bg = COLOR_LIGHT_GRAY if row_idx % 2 == 0 else COLOR_WHITE
            table.cell(row_idx, 0).text = feature.name
            _style_cell(table.cell(row_idx, 0), Pt(9), bold=True, bg=bg, fg=COLOR_TEXT)
            table.cell(row_idx, 1).text = val_a
            _style_cell(table.cell(row_idx, 1), Pt(8), bg=bg, fg=COLOR_TEXT)
            table.cell(row_idx, 2).text = val_b
            _style_cell(table.cell(row_idx, 2), Pt(8), bg=bg, fg=COLOR_TEXT)


class CaseStudySlideBuilder(SlideBuilder):
    """Builds the case studies slide."""

    def build(self, slide, result: ComparisonResult) -> None:
        if not result.case_studies:
            return

        # Title
        txBox = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.25), CONTENT_WIDTH, Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{result.product_a.name} Case Studies — {result.industry.name}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        y = Inches(0.85)
        for cs in result.case_studies[:3]:
            # Title + customer
            txTitle = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.35))
            tf_t = txTitle.text_frame
            p_t = tf_t.paragraphs[0]
            p_t.text = f"{cs.title}  —  {cs.customer}"
            p_t.font.size = Pt(12)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_ACCENT

            # Summary
            txS = slide.shapes.add_textbox(CONTENT_LEFT, y + Inches(0.38), CONTENT_WIDTH, Inches(0.5))
            tf_s = txS.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.text = cs.summary
            p_s.font.size = Pt(9)
            p_s.font.color.rgb = COLOR_TEXT

            # Results
            if cs.results:
                txR = slide.shapes.add_textbox(CONTENT_LEFT, y + Inches(0.9), CONTENT_WIDTH, Inches(0.4))
                tf_r = txR.text_frame
                tf_r.word_wrap = True
                p_r = tf_r.paragraphs[0]
                p_r.text = f"Results: {cs.results}"
                p_r.font.size = Pt(8)
                p_r.font.italic = True
                p_r.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)

            y += Inches(1.5)


class LLMSuggestionSlideBuilder(SlideBuilder):
    """Builds the LLM-generated suggestions slide."""

    def build(self, slide, result: ComparisonResult) -> None:
        # Title
        txBox = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.25), CONTENT_WIDTH, Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Expert Analysis & Recommendations"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        # LLM Summary
        summary_text = result.llm_summary or "No analysis available."
        txBox2 = slide.shapes.add_textbox(CONTENT_LEFT, Inches(0.85), Inches(9.0), Inches(3.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, line in enumerate(summary_text.split("\n")):
            para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            para.text = line
            para.font.size = Pt(9)
            para.font.color.rgb = COLOR_TEXT
            para.space_after = Pt(3)

        # Expert quotes
        if result.expert_advice_list:
            y_pos = Inches(4.4)
            txBox3 = slide.shapes.add_textbox(CONTENT_LEFT, y_pos, CONTENT_WIDTH, Inches(0.3))
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = "Field Expert Notes"
            p3.font.size = Pt(11)
            p3.font.bold = True
            p3.font.color.rgb = COLOR_ACCENT

            txBox4 = slide.shapes.add_textbox(CONTENT_LEFT, y_pos + Inches(0.35), Inches(9.0), Inches(0.8))
            tf4 = txBox4.text_frame
            tf4.word_wrap = True
            for i, advice in enumerate(result.expert_advice_list[:2]):
                para = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
                para.text = f'"{advice.advice_text[:120]}..." — {advice.author}'
                para.font.size = Pt(8)
                para.font.italic = True
                para.font.color.rgb = COLOR_TEXT
                para.space_after = Pt(4)


def _style_cell(cell, font_size, bold=False, bg=None, fg=None):
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        if fg:
            paragraph.font.color.rgb = fg
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


class PPTGeneratorFactory:
    """Factory for creating the slide builder pipeline."""

    @staticmethod
    def create_builders() -> List[SlideBuilder]:
        return [
            FeatureTableSlideBuilder(),
            CaseStudySlideBuilder(),
            LLMSuggestionSlideBuilder(),
        ]


class PPTGenerator:
    """Generates Battle Card Flash PPT files using the template."""

    def __init__(self):
        self._builders = PPTGeneratorFactory.create_builders()

    def generate(self, result: ComparisonResult, output_path: str) -> str:
        prs = Presentation(TEMPLATE_PATH)

        # Template structure: slide 0 = title, slide 1 = content template, slide 2 = thanks
        title_slide = prs.slides[0]
        content_template = prs.slides[1]

        # Add title text to the title slide
        txBox = title_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{result.product_a.name} vs. {result.product_b.name}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = result.industry.name
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        p2.alignment = PP_ALIGN.CENTER

        # Clone the content template for each builder and populate it
        for builder in self._builders:
            new_slide = _clone_slide(prs, content_template)
            builder.build(new_slide, result)

        # Move the thanks slide (originally index 2, now still at index 2) to the end
        # by removing and re-appending it.
        # Current order: [title, content_tmpl, thanks, feat, case, llm]
        # Desired order: [title, feat, case, llm, thanks]
        slide_list = prs.slides._sldIdLst
        slide_ids = list(slide_list)

        # Remove original content template (index 1) and thanks (index 2)
        tmpl_id = slide_ids[1]
        thanks_id = slide_ids[2]
        slide_list.remove(tmpl_id)
        slide_list.remove(thanks_id)
        # Re-append thanks at the end
        slide_list.append(thanks_id)

        prs.save(output_path)
        return output_path
